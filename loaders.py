"""
문서 로더 모듈 - PDF, Markdown, Text 파일 로딩 및 전처리
"""
import os
import unicodedata
import re
from typing import List, Dict, Any
from pathlib import Path
import logging

# LangChain document loaders
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader, TextLoader

# Office Loaders
try:
    import docx
    import pptx
    import openpyxl
except ImportError:
    pass

# OCR Loaders
try:
    import pytesseract
    from pdf2image import convert_from_path
    from PIL import Image
except ImportError:
    pass

# 로깅 설정
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentLoader:
    """다양한 형식의 문서를 로드하고 전처리하는 클래스"""
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.md', '.txt', '.docx', '.pptx', '.xlsx'}
    
    def __init__(self):
        self.stats = {
            'total_files': 0,
            'success_files': 0,
            'failed_files': 0,
            'error_messages': []
        }
    
    def load_file(self, file_path: str) -> List[Document]:
        """
        파일 확장자에 따라 적절한 로더를 선택하여 문서 로드
        
        Args:
            file_path: 로드할 파일 경로
            
        Returns:
            Document 객체 리스트
        """
        path = Path(file_path)
        
        if not path.exists():
            error_msg = f"파일을 찾을 수 없습니다: {file_path}"
            logger.error(error_msg)
            self.stats['failed_files'] += 1
            self.stats['error_messages'].append(error_msg)
            return []
        
        extension = path.suffix.lower()
        
        if extension not in self.SUPPORTED_EXTENSIONS:
            error_msg = f"지원하지 않는 파일 형식입니다: {extension}"
            logger.error(error_msg)
            self.stats['failed_files'] += 1
            self.stats['error_messages'].append(error_msg)
            return []
        
        self.stats['total_files'] += 1
        
        try:
            if extension == '.pdf':
                documents = self._load_pdf(file_path)
            elif extension == '.md':
                documents = self._load_markdown(file_path)
            elif extension == '.txt':
                documents = self._load_text(file_path)
            elif extension == '.docx':
                documents = self._load_docx(file_path)
            elif extension == '.pptx':
                documents = self._load_pptx(file_path)
            elif extension == '.xlsx':
                documents = self._load_xlsx(file_path)
            else:
                documents = []
            
            # 전처리 적용
            documents = self._preprocess_documents(documents, file_path)
            
            self.stats['success_files'] += 1
            logger.info(f"성공적으로 로드: {file_path} ({len(documents)} 페이지)")
            
            return documents
            
        except Exception as e:
            error_msg = f"파일 로드 실패 ({file_path}): {str(e)}"
            logger.error(error_msg)
            self.stats['failed_files'] += 1
            self.stats['error_messages'].append(error_msg)
            return []
    
    def _load_pdf(self, file_path: str) -> List[Document]:
        """PDF 파일 로드 (OCR Fallback 포함)"""
        try:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            
            # 텍스트 추출 품질 확인 및 OCR Fallback
            ocr_needed = False
            total_text_len = sum(len(doc.page_content.strip()) for doc in documents)
            avg_text_len = total_text_len / len(documents) if documents else 0
            
            # 텍스트가 너무 적으면(페이지당 50자 미만) 스캔본으로 간주
            if avg_text_len < 50:
                logger.info(f"텍스트 추출 품질 저조 (평균 {avg_text_len:.1f}자). OCR을 시도합니다: {file_path}")
                ocr_needed = True
            
            if ocr_needed:
                return self._ocr_pdf(file_path)
                
            return documents
        except Exception as e:
            logger.error(f"PDF 로드 오류: {str(e)}")
            raise

    def _ocr_pdf(self, file_path: str) -> List[Document]:
        """PDF OCR 처리"""
        try:
            logger.info("PDF를 이미지로 변환 중...")
            images = convert_from_path(file_path)
            documents = []
            
            for i, image in enumerate(images):
                logger.info(f"OCR 처리 중 (페이지 {i+1}/{len(images)})...")
                # 한글+영어 OCR
                text = pytesseract.image_to_string(image, lang='kor+eng')
                
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={"source": file_path, "page": i+1, "ocr": True}
                    ))
            
            return documents
        except Exception as e:
            logger.error(f"OCR 처리 실패: {str(e)}")
            # OCR 실패 시 빈 리스트 반환보다는 에러 전파가 나을 수 있음, 
            # 하지만 여기서는 로더 특성상 빈 리스트나 부분 성공도 고려
            if 'poppler' in str(e).lower():
                logger.error("Poppler가 설치되지 않은 것 같습니다. 'brew install poppler'를 확인하세요.")
            raise
    
    def _load_markdown(self, file_path: str) -> List[Document]:
        """Markdown 파일 로드"""
        try:
            # UTF-8 인코딩 시도
            return self._load_text_with_encoding(file_path, 'utf-8')
        except UnicodeDecodeError:
            # CP949 인코딩 시도 (한글 Windows 환경)
            try:
                logger.warning(f"UTF-8 디코딩 실패, CP949로 재시도: {file_path}")
                return self._load_text_with_encoding(file_path, 'cp949')
            except Exception as e:
                logger.error(f"Markdown 로드 오류: {str(e)}")
                raise
    
    def _load_text(self, file_path: str) -> List[Document]:
        """Text 파일 로드"""
        try:
            # UTF-8 인코딩 시도
            return self._load_text_with_encoding(file_path, 'utf-8')
        except UnicodeDecodeError:
            # CP949 인코딩 시도
            try:
                logger.warning(f"UTF-8 디코딩 실패, CP949로 재시도: {file_path}")
                return self._load_text_with_encoding(file_path, 'cp949')
            except Exception as e:
                logger.error(f"Text 로드 오류: {str(e)}")
                raise
    
    def _load_text_with_encoding(self, file_path: str, encoding: str) -> List[Document]:
        """특정 인코딩으로 텍스트 파일 로드"""
        try:
            loader = TextLoader(file_path, encoding=encoding)
            documents = loader.load()
            return documents
        except Exception as e:
            logger.error(f"텍스트 로드 오류 ({encoding}): {str(e)}")
            raise

    def _load_docx(self, file_path: str) -> List[Document]:
        """Word 파일 로드"""
        try:
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            text = "\n\n".join(full_text)
            return [Document(page_content=text, metadata={"source": file_path})]
        except Exception as e:
            logger.error(f"Word 로드 오류: {str(e)}")
            raise

    def _load_pptx(self, file_path: str) -> List[Document]:
        """PowerPoint 파일 로드"""
        try:
            prs = pptx.Presentation(file_path)
            documents = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text)
                
                if slide_text:
                    text = "\n".join(slide_text)
                    documents.append(Document(
                        page_content=text, 
                        metadata={"source": file_path, "page": i+1}
                    ))
            
            return documents
        except Exception as e:
            logger.error(f"PPTX 로드 오류: {str(e)}")
            raise

    def _load_xlsx(self, file_path: str) -> List[Document]:
        """Excel 파일 로드"""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            documents = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = []
                
                # 헤더 읽기
                headers = []
                for row in sheet.iter_rows(min_row=1, max_row=1, values_only=True):
                    headers = [str(h) if h is not None else f"Col{i}" for i, h in enumerate(row)]
                
                # 데이터 읽기
                for row in sheet.iter_rows(min_row=2, values_only=True):
                    row_text = []
                    has_data = False
                    for i, cell in enumerate(row):
                        if cell is not None:
                            header = headers[i] if i < len(headers) else f"Col{i}"
                            row_text.append(f"{header}: {cell}")
                            has_data = True
                    
                    if has_data:
                        sheet_text.append(", ".join(row_text))
                
                if sheet_text:
                    text = f"Sheet: {sheet_name}\n" + "\n".join(sheet_text)
                    documents.append(Document(
                        page_content=text,
                        metadata={"source": file_path, "sheet": sheet_name}
                    ))
            
            return documents
        except Exception as e:
            logger.error(f"Excel 로드 오류: {str(e)}")
            raise
    
    def _preprocess_documents(self, documents: List[Document], file_path: str) -> List[Document]:
        """
        문서 전처리 적용
        - 공백 정규화
        - 한글 자모 복구
        - 빈 문서 제거
        """
        processed_docs = []
        
        for i, doc in enumerate(documents):
            # 빈 문서 스킵
            if not doc.page_content or doc.page_content.strip() == "":
                logger.warning(f"빈 문서 발견 (페이지 {i+1}), 건너뜀")
                continue
            
            # 전처리 적용
            cleaned_text = self._clean_text(doc.page_content)
            
            # 메타데이터 보강
            doc.page_content = cleaned_text
            doc.metadata['file_path'] = file_path
            doc.metadata['file_name'] = Path(file_path).name
            
            # 페이지 번호가 없으면 추가
            if 'page' not in doc.metadata:
                doc.metadata['page'] = i + 1
            
            processed_docs.append(doc)
        
        return processed_docs
    
    def _clean_text(self, text: str) -> str:
        """
        텍스트 정리 및 정규화
        - 공백 정규화
        - 한글 자모 복구 (NFC 정규화)
        - 과도한 줄바꿈 제거
        """
        # 한글 자모 복구 (유니코드 정규화)
        text = unicodedata.normalize('NFC', text)
        
        # 다중 공백을 단일 공백으로 치환
        text = re.sub(r'[ \t]+', ' ', text)
        
        # 3개 이상의 연속된 줄바꿈을 2개로 치환
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # 앞뒤 공백 제거
        text = text.strip()
        
        return text
    
    def get_stats(self) -> Dict[str, Any]:
        """로딩 통계 반환"""
        return self.stats
    
    def reset_stats(self):
        """통계 초기화"""
        self.stats = {
            'total_files': 0,
            'success_files': 0,
            'failed_files': 0,
            'error_messages': []
        }


def load_documents(file_paths: List[str]) -> List[Document]:
    """
    여러 파일을 한 번에 로드하는 헬퍼 함수
    
    Args:
        file_paths: 로드할 파일 경로 리스트
        
    Returns:
        모든 Document 객체를 합친 리스트
    """
    loader = DocumentLoader()
    all_documents = []
    
    for file_path in file_paths:
        documents = loader.load_file(file_path)
        all_documents.extend(documents)
    
    # 통계 출력
    stats = loader.get_stats()
    logger.info(f"=== 문서 로딩 완료 ===")
    logger.info(f"총 파일: {stats['total_files']}")
    logger.info(f"성공: {stats['success_files']}")
    logger.info(f"실패: {stats['failed_files']}")
    
    if stats['error_messages']:
        logger.warning(f"오류 메시지:")
        for msg in stats['error_messages']:
            logger.warning(f"  - {msg}")
    
    return all_documents


if __name__ == "__main__":
    # 테스트 코드
    import sys
    
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        docs = load_documents([test_file])
        
        print(f"\n로드된 문서 수: {len(docs)}")
        for i, doc in enumerate(docs[:3]):  # 처음 3개만 출력
            print(f"\n--- Document {i+1} ---")
            print(f"Length: {len(doc.page_content)} chars")
            print(f"Metadata: {doc.metadata}")
            print(f"Preview: {doc.page_content[:200]}...")
    else:
        print("사용법: python loaders.py <file_path>")
