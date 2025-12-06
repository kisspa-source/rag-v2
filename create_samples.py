import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
import os
from pathlib import Path

# Ensure data directory exists
os.makedirs("data/samples", exist_ok=True)

def create_sample_docx(filename="data/samples/sample.docx"):
    doc = docx.Document()
    
    # Title
    doc.add_heading('Optimizing Local RAG Systems: macOS vs Linux Performance Analysis', 0)
    
    # Abstract
    doc.add_heading('Abstract (초록)', level=1)
    doc.add_paragraph(
        "본 문서는 로컬 환경(On-premise)에서 구축된 RAG(Retrieval-Augmented Generation) 시스템의 성능을 "
        "macOS(Apple Silicon)와 Linux(NVIDIA GPU) 환경에서 비교 분석합니다. "
        "We investigate the impact of quantization, inference engine optimization (Ollama, vLLM), "
        "and retrieval latencies. 실험 결과, Linux 환경의 CUDA 가속이 처리량(Throughput) 면에서 우세했으나, "
        "macOS의 통합 메모리 구조(Unified Memory Architecture)가 대규모 모델 로딩에 있어 비용 효율적인 대안임을 확인했습니다."
    )
    
    # Introduction
    doc.add_heading('1. Introduction', level=1)
    doc.add_paragraph(
        "Recent advancements in Large Language Models (LLMs) have democratized access to powerful AI capabilities. "
        "However, data privacy concerns have driven the adoption of Local RAG systems. "
        "기업 내부망이나 개인 PC에서 외부 API 호출 없이 문서를 검색하고 답변을 생성하는 시스템의 수요가 급증하고 있습니다."
    )
    doc.add_paragraph(
        "특히 Apple의 M-series 칩셋(M1, M2, M3)이 NPU(Neural Engine)를 탑재하면서 "
        "개인용 노트북에서도 7B~13B 파라미터 규모의 모델을 원활하게 구동할 수 있게 되었습니다. "
        "On the other hand, Linux workstations with dedicated NVIDIA GPUs remain the gold standard for low-latency inference."
    )
    
    # Methodology
    doc.add_heading('2. System Architecture', level=1)
    doc.add_paragraph("본 실험에 사용된 RAG 파이프라인의 구성요소는 다음과 같습니다:")
    
    p = doc.add_paragraph(style='List Bullet')
    p.add_run('Embeddings: ').bold = True
    p.add_run('BAAI/bge-m3 (Multi-lingual, High performance)')
    
    p = doc.add_paragraph(style='List Bullet')
    p.add_run('Vector DB: ').bold = True
    p.add_run('ChromaDB (Local persistent storage)')
    
    p = doc.add_paragraph(style='List Bullet')
    p.add_run('LLM Inference: ').bold = True
    p.add_run('Ollama (Llama.cpp backend) running qwen2:7b and llama3:8b')
    
    p = doc.add_paragraph(style='List Bullet')
    p.add_run('Reranker: ').bold = True
    p.add_run('BAAI/bge-reranker-v2-m3 (Cross-encoder for higher accuracy)')

    # Results
    doc.add_heading('3. Performance Benchmarks', level=1)
    doc.add_paragraph(
        "We measured 'Time to First Token (TTFT)' and 'Tokens Per Second (TPS)'. "
        "Linux 환경(RTX 4090)에서는 TPS가 80+에 달했으나, macOS M2 Max에서도 Metal 가속을 통해 40+ TPS의 준수한 성능을 보였습니다. "
        "특히 Reranking 단계에서는 CPU 의존도가 높아, 멀티코어 성능이 중요한 변수로 작용했습니다."
    )
    
    # Conclusion
    doc.add_heading('4. Conclusion', level=1)
    doc.add_paragraph(
        "For production-grade local RAG, Linux with CUDA is recommended. "
        "하지만 개발 및 소규모 배포 환경에서는 macOS가 전력 효율성과 편의성 면에서 압도적인 장점을 가집니다. "
        "Future work will explore distributed inference across heterogeneous devices."
    )

    doc.save(filename)
    print(f"Created {filename}")

def create_sample_pptx(filename="data/samples/sample.pptx"):
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Local RAG Implementation Strategy"
    slide.placeholders[1].text = "Building Secure, Private AI Search\n\nComparing Linux & macOS Infrastructures"
    
    # Slide 2: Why Local RAG?
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Why Local RAG?"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Privacy & Security (보안 및 개인정보 보호)"
    p = tf.add_paragraph()
    p.text = "Sensitive data never leaves the premise (데이터 외부 유출 방지)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero-trust environments support"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cost Efficiency (비용 효율성)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "No recurring API costs (Token throughput cost optimized)"
    p.level = 1
    
    # Slide 3: macOS vs Linux
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Platform Comparison: macOS vs Linux"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "macOS (Apple Silicon)"
    p = tf.add_paragraph()
    p.text = "Unified Memory allows loading massive models (up to 128GB+ VRAM equivalent)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Great for development, high energy efficiency"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Linux (NVIDIA / AMD)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "CUDA optimization provides best-in-class raw performance"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Better ecosystem support (vLLM, TGI, TensorRT-LLM)"
    p.level = 1

    prs.save(filename)
    print(f"Created {filename}")

def create_sample_xlsx(filename="data/samples/sample.xlsx"):
    wb = Workbook()
    ws = wb.active
    ws.title = "Benchmark Results"
    
    # Headers
    headers = ["Model", "Platform", "Hardware", "Quantization", "Prompt Tokens", "Context Length", "TTFT (ms)", "TPS (Tokens/sec)", "Memory (GB)"]
    ws.append(headers)
    
    # Data
    data = [
        ["Llama-3-8B", "Linux", "NVIDIA RTX 4090", "Q4_K_M", 128, 4096, 12, 115.4, 6.2],
        ["Llama-3-8B", "Linux", "NVIDIA RTX 3090", "Q4_K_M", 128, 4096, 15, 85.2, 6.2],
        ["Llama-3-8B", "macOS", "M3 Max (30-core GPU)", "Q4_0", 128, 4096, 45, 52.8, 5.8],
        ["Llama-3-8B", "macOS", "M1 Pro", "Q4_0", 128, 4096, 120, 18.5, 5.8],
        ["Qwen2-72B", "Linux", "2x A100 80GB", "FP16", 512, 8192, 180, 22.1, 144.5],
        ["Qwen2-72B", "macOS", "M2 Ultra (192GB RAM)", "Q4_K_M", 512, 8192, 450, 8.4, 48.2],
        ["Mistral-Large", "Linux", "H100 Hopper", "FP8", 256, 32000, 8, 140.2, 24.5],
    ]
    
    for row in data:
        ws.append(row)
        
    # Another sheet for Retrieval Latency
    ws2 = wb.create_sheet("Retrieval Latency")
    ws2.append(["Algorithm", "Query Type", "Index Size", "Latency (ms)"])
    ws2.append(["BM25", "Keyword", "100k docs", 15])
    ws2.append(["HNSW (Vector)", "Semantic", "100k docs", 45])
    ws2.append(["Hybrid (BM25+Vector)", "Complex", "100k docs", 58])
    ws2.append(["RRF Ranking", "Re-ranking", "Top 50 docs", 120])
    
    wb.save(filename)
    print(f"Created {filename}")

def create_sample_md(filename="data/samples/sample.md"):
    content = """# Local RAG Setup Guide (로컬 RAG 구축 가이드)

## 1. Environment Setup (환경 설정)

To set up a robust local RAG (Retrieval-Augmented Generation) system, we recommend using **Ollama** for model inference and **LangChain** for orchestration.

### 1.1 Prerequisites
*   OS: macOS (Ventura or later) or Linux (Ubuntu 22.04 LTS)
*   Python 3.10+
*   RAM: Min 16GB (32GB+ recommended for 70B models)

### 1.2 Installation

```bash
# 1. Install Ollama
curl -fsSL https://ollama.com/install.sh | sh

# 2. Pull Embedding & LLM Models
ollama pull qwen2:7b
ollama pull mxbai-embed-large

# 3. Create Python Virtual Environment
python3 -m venv venv
source venv/bin/activate
pip install langchain langchain-community chromadb
```

## 2. RAG Pipeline Configuration

로컬 RAG 파이프라인은 보안성이 뛰어나지만, 검색 정확도(Retrieval Accuracy)를 높이기 위해 정교한 튜닝이 필요합니다.

### 2.1 Chunking Strategy
Document splitting is crucial.
*   **Fixed-size chunking**: Simple but may break context.
*   **RecursiveCharacterTextSplitter**: Respects sentence structure (Recommended).
*   **Semantic Chunking**: Splits based on meaning changes (Advanced).

> **Tip**: 한글 문서의 경우 `chunk_size=500~800`, `chunk_overlap=50~100` 정도가 적당합니다.

### 2.2 Vector Database
For local lightweight setups, **ChromaDB** or **FAISS** are popular.
ChromaDB saves raw data and embeddings in `./chroma_db`, making it persistent across restarts.

## 3. TroubleShooting (문제 해결)

| Issue | Possible Cause | Solution |
|-------|----------------|----------|
| Low TPS | CPU bottleneck | Enable GPU/Metal acceleration in Ollama |
| Hallucination | Poor Retrieval | Use Hybrid Search (BM25 + Vector) |
| OOM (Out of Memory) | Context window too large | Reduce `context_count` or quantized model |

---
*Created by RAG Automation Script*
"""
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"Created {filename}")

def create_sample_txt(filename="data/samples/sample.txt"):
    content = """
2024-05-20 10:00:01 [INFO] RAG_System: Starting initialization sequence...
2024-05-20 10:00:02 [INFO] Loader: Loaded configuration from config/config.yaml
2024-05-20 10:00:03 [INFO] GPU_Check: NVIDIA CUDA detected (Count: 2, Model: A100-80GB)
2024-05-20 10:00:03 [WARN] GPU_Check: Driver version 535.122 is slightly outdated. Recommend update to 550+.
2024-05-20 10:00:05 [INFO] Model_Load: Loading 'llama3:70b-instruct-q4_k_m'...
2024-05-20 10:00:15 [INFO] Model_Load: Successfully loaded to VRAM (Offload layers: 83/83)
2024-05-20 10:00:16 [INFO] VectorDB: Connecting to ChromaDB (Persistent) at ./chroma_db
2024-05-20 10:00:17 [INFO] Indexer: Found 12,450 documents in index.
2024-05-20 10:00:17 [INFO] Server: API Server running on http://0.0.0.0:8000

--- REQUEST LOG ---
2024-05-20 10:05:23 [REQ] User_ID: user_123 | Query: "리눅스에서 로컬 RAG 성능 최적화 방법은?"
2024-05-20 10:05:23 [SEARCH] Hybrid Search triggered (BM25 + Vector)
2024-05-20 10:05:24 [RETRIEVAL] Found 50 candidates. Top source: 'linux_optimization_guide.pdf' (Score: 0.89)
2024-05-20 10:05:24 [RERANK] BAAI/bge-reranker-v2-m3 re-ordered top 5 candidates.
2024-05-20 10:05:24 [GEN] Context len: 2450 tokens. Generating answer...
2024-05-20 10:05:28 [RES] Sent 450 tokens. (Speed: 112.5 t/s). Latency: 4.2s.

--- SYSTEM METRICS ---
CPU Usage: 12%
RAM Usage: 24.5 GB / 128 GB
VRAM Usage: 78.2 GB / 160 GB (Dual GPU)
Temperature: GPU0 65C, GPU1 68C
Power Draw: 550W (Peak)
"""
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"Created {filename}")

if __name__ == "__main__":
    create_sample_docx()
    create_sample_pptx()
    create_sample_xlsx()
    create_sample_md()
    create_sample_txt()

